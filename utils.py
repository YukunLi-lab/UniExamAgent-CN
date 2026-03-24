"""
UniExamAgent-CN 工具模块
文件处理、爬虫、PPT 提取等工具函数
"""

import os
import re
import time
import json
import hashlib
from pathlib import Path
from typing import List, Dict, Any, Optional
from urllib.parse import urljoin, urlparse

import requests
from bs4 import BeautifulSoup
from loguru import logger
from tqdm import tqdm

from ocr_corrector import correct_ocr_formulas

# ==================== 文件处理 ====================
class FileProcessor:
    """多格式文件处理器"""

    SUPPORTED_FORMATS = [".pptx", ".pdf", ".txt", ".jpg", ".png", ".jpeg"]

    @staticmethod
    def process_file(file_path: Path) -> str:
        """根据文件类型调用对应处理器"""
        suffix = file_path.suffix.lower()

        processors = {
            ".pptx": FileProcessor._extract_from_pptx,
            ".pdf": FileProcessor._extract_from_pdf,
            ".txt": FileProcessor._extract_from_txt,
            ".jpg": FileProcessor._extract_from_image,
            ".png": FileProcessor._extract_from_image,
            ".jpeg": FileProcessor._extract_from_image,
        }

        processor = processors.get(suffix)
        if not processor:
            logger.warning(f"不支持的文件格式: {suffix}")
            return ""

        try:
            text = processor(file_path)
            # 对图片/PDF 的 OCR 结果进行公式纠错
            if suffix in [".jpg", ".png", ".jpeg", ".pdf"]:
                text = correct_ocr_formulas(text, wrap=True)
            return text
        except Exception as e:
            logger.error(f"处理文件失败 {file_path}: {e}")
            return ""

    @staticmethod
    def _extract_from_pptx(pptx_path: Path) -> str:
        """从 PPTX 提取文本"""
        from pptx import Presentation

        texts = []
        prs = Presentation(str(pptx_path))

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"[幻灯片 {slide_num}]"]

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

                # 处理表格
                if shape.has_table:
                    table_text = []
                    for row in shape.table.rows:
                        row_text = [cell.text for cell in row.cells]
                        table_text.append(" | ".join(row_text))
                    if table_text:
                        slide_text.append(" | ".join(table_text))

            texts.append("\n".join(slide_text))

        logger.info(f"从 PPTX 提取 {len(prs.slides)} 张幻灯片")
        return "\n\n".join(texts)

    @staticmethod
    def _extract_from_pdf(pdf_path: Path) -> str:
        """从 PDF 提取文本（含 OCR 图片化页面）"""
        import pdfplumber
        import pytesseract
        from PIL import Image
        import pypdfium2 as pdfium

        texts = []

        # 先尝试用 pdfplumber 提取文字
        with pdfplumber.open(pdf_path) as pdf:
            page_count = len(pdf.pages)

            for page_num, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text and len(text.strip()) > 10:
                    texts.append(f"[第 {page_num} 页]\n{text}")
                else:
                    # 页面文字少或为空，可能是扫描版，进行 OCR
                    logger.info(f"第 {page_num} 页文字少，使用 OCR 提取...")
                    ocr_text = FileProcessor._ocr_pdf_page(pdf_path, page_num - 1)
                    if ocr_text:
                        texts.append(f"[第 {page_num} 页 - OCR]\n{ocr_text}")

        # 如果 pdfplumber 提取的文字太少，尝试全页 OCR
        total_text = "\n\n".join(texts)
        if len(total_text.strip()) < 100:
            logger.info("pdfplumber 提取文字过少，尝试全页 OCR...")
            full_ocr_text = FileProcessor._ocr_entire_pdf(pdf_path)
            if full_ocr_text:
                texts = [f"[全文档 OCR]\n{full_ocr_text}"]

        logger.info(f"从 PDF 提取 {page_count} 页")
        return "\n\n".join(texts)

    @staticmethod
    def _ocr_pdf_page(pdf_path: Path, page_index: int) -> str:
        """对 PDF 单页进行 OCR"""
        import pypdfium2 as pdfium
        import pytesseract
        import easyocr
        import numpy as np
        from PIL import Image
        import io

        try:
            pdf = pdfium.PdfDocument(str(pdf_path))
            page = pdf[page_index]
            pil_image = page.render(scale=2).to_pil()
            img_array = np.array(pil_image)

            # 优先尝试 Tesseract OCR
            try:
                text = pytesseract.image_to_string(pil_image, lang="chi_sim+eng")
                if text and len(text.strip()) > 5:
                    logger.info(f"页面 {page_index + 1} OCR (Tesseract) 提取 {len(text)} 字符")
                    return text
            except Exception as e:
                logger.warning(f"Tesseract OCR 失败: {e}")

            # 备用 EasyOCR
            try:
                reader = easyocr.Reader(["ch_sim", "en"], verbose=False)
                results = reader.readtext(img_array)
                text = "\n".join([r[1] for r in results if r[2] > 0.3])
                if text and len(text.strip()) > 5:
                    logger.info(f"页面 {page_index + 1} OCR (EasyOCR) 提取 {len(text)} 字符")
                    return text
            except Exception as e:
                logger.warning(f"EasyOCR 失败: {e}")

        except Exception as e:
            logger.warning(f"页面 OCR 失败: {e}")
        return ""

    @staticmethod
    def _ocr_entire_pdf(pdf_path: Path) -> str:
        """对整个 PDF 进行 OCR（备选方案）"""
        import pypdfium2 as pdfium
        import easyocr
        import numpy as np
        from PIL import Image

        try:
            pdf = pdfium.PdfDocument(str(pdf_path))
            all_texts = []

            # 使用 EasyOCR（支持中文+英文）
            reader = easyocr.Reader(["ch_sim", "en"], verbose=False)

            for page_num in range(len(pdf)):
                page = pdf[page_num]
                pil_image = page.render(scale=2).to_pil()
                img_array = np.array(pil_image)

                try:
                    results = reader.readtext(img_array)
                    text = "\n".join([r[1] for r in results if r[2] > 0.3])
                    if text and len(text.strip()) > 5:
                        all_texts.append(f"[第 {page_num + 1} 页]\n{text}")
                except Exception as e:
                    logger.warning(f"页面 {page_num + 1} OCR 失败: {e}")

            if all_texts:
                logger.info(f"全文档 OCR 提取 {len(pdf)} 页，共 {len(all_texts)} 页有文字")
                return "\n\n".join(all_texts)
        except Exception as e:
            logger.warning(f"全文档 OCR 失败: {e}")
        return ""

    @staticmethod
    def _extract_from_txt(txt_path: Path) -> str:
        """从 TXT 读取文本"""
        with open(txt_path, "r", encoding="utf-8") as f:
            return f.read()

    @staticmethod
    def _extract_from_image(image_path: Path) -> str:
        """从图片 OCR 提取文本"""
        from PIL import Image

        # 优先尝试 Tesseract OCR
        try:
            import pytesseract
            img = Image.open(image_path)
            text = pytesseract.image_to_string(img, lang="chi_sim+eng")
            if text and len(text.strip()) > 5:
                logger.info(f"图片 OCR (Tesseract) 提取 {len(text)} 字符")
                return text
        except Exception as e:
            logger.warning(f"Tesseract OCR 失败: {e}")

        # 备用 EasyOCR
        try:
            import easyocr
            reader = easyocr.Reader(["ch_sim", "en"], verbose=False)
            results = reader.readtext(str(image_path))
            text = "\n".join([r[1] for r in results if r[2] > 0.3])
            if text and len(text.strip()) > 5:
                logger.info(f"图片 OCR (EasyOCR) 提取 {len(text)} 字符")
                return text
        except Exception as e:
            logger.warning(f"EasyOCR 失败: {e}")

        return ""

    @staticmethod
    def process_uploaded_files(upload_dir: Path) -> List[Dict[str, Any]]:
        """批量处理上传目录中的所有文件"""
        results = []

        for file_path in tqdm(list(upload_dir.glob("*")), desc="处理文件"):
            if file_path.suffix.lower() not in FileProcessor.SUPPORTED_FORMATS:
                continue

            content = FileProcessor.process_file(file_path)
            if content:
                results.append({
                    "file_name": file_path.name,
                    "file_type": file_path.suffix.lower(),
                    "content": content,
                    "char_count": len(content),
                    "file_hash": hashlib.md5(content.encode()).hexdigest()
                })

        logger.info(f"共处理 {len(results)} 个文件")
        return results


# ==================== 网络爬虫 ====================
class CourseCrawler:
    """课程资料爬虫"""

    def __init__(self, timeout: int = 30, delay: float = 1.0):
        self.timeout = timeout
        self.delay = delay
        self.session = requests.Session()
        self.session.headers.update({
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
        })

    def crawl(self, university: str, course_name: str,
              keywords: List[str], max_results: int = 50) -> List[Dict]:
        """爬取课程相关资料"""
        results = []

        search_engines = [
            ("mooc", self._crawl_mooc),
            ("zhihu", self._crawl_zhihu),
            ("baidu_wenku", self._crawl_baidu_wenku),
        ]

        for engine_name, crawl_func in search_engines:
            logger.info(f"正在爬取 {engine_name}...")

            try:
                items = crawl_func(university, course_name, keywords, max_results)
                results.extend(items)
                time.sleep(self.delay)
            except Exception as e:
                logger.error(f"爬取 {engine_name} 失败: {e}")

        return results[:max_results]

    def _crawl_mooc(self, university: str, course_name: str,
                    keywords: List[str], max_results: int) -> List[Dict]:
        """爬取中国大学 MOOC"""
        results = []

        search_url = f"https://www.icourse163.org/search.htm"
        params = {
            "searchKey": f"{university} {course_name}",
            "type": "course",
            "orgId": ""
        }

        try:
            response = self.session.get(search_url, params=params, timeout=self.timeout)
            soup = BeautifulSoup(response.text, "lxml")

            # 提取课程卡片
            course_cards = soup.select(".course-card")
            for card in course_cards[:max_results]:
                title = card.select_one(".course-name")
                if title:
                    results.append({
                        "source_url": "https://www.icourse163.org",
                        "source_name": "中国大学 MOOC",
                        "title": title.text.strip(),
                        "content_summary": "MOOC 课程视频与课件",
                        "relevance_score": 0.9,
                        "crawl_status": "success"
                    })
        except Exception as e:
            logger.warning(f"MOOC 爬取失败: {e}")

        return results

    def _crawl_zhihu(self, university: str, course_name: str,
                     keywords: List[str], max_results: int) -> List[Dict]:
        """爬取知乎"""
        results = []

        search_url = f"https://www.zhihu.com/search"
        params = {"type": "content", "q": f"{university} {course_name}"}

        try:
            response = self.session.get(search_url, params=params, timeout=self.timeout)
            soup = BeautifulSoup(response.text, "lxml")

            articles = soup.select(".List-item")
            for article in articles[:max_results]:
                title_elem = article.select_one(".ContentItem-title")
                if title_elem:
                    results.append({
                        "source_url": "https://www.zhihu.com",
                        "source_name": "知乎",
                        "title": title_elem.text.strip(),
                        "content_summary": "知乎问答与文章",
                        "relevance_score": 0.7,
                        "crawl_status": "success"
                    })
        except Exception as e:
            logger.warning(f"知乎爬取失败: {e}")

        return results

    def _crawl_baidu_wenku(self, university: str, course_name: str,
                          keywords: List[str], max_results: int) -> List[Dict]:
        """爬取百度文库"""
        results = []

        search_url = "https://wenku.baidu.com/search"
        params = {"word": f"{university} {course_name}"}

        try:
            response = self.session.get(search_url, params=params, timeout=self.timeout)
            soup = BeautifulSoup(response.text, "lxml")

            doc_cards = soup.select(".search-result-item")
            for card in doc_cards[:max_results]:
                title = card.select_one(".title")
                if title:
                    results.append({
                        "source_url": "https://wenku.baidu.com",
                        "source_name": "百度文库",
                        "title": title.text.strip(),
                        "content_summary": "文库文档与资料",
                        "relevance_score": 0.8,
                        "crawl_status": "success"
                    })
        except Exception as e:
            logger.warning(f"百度文库爬取失败: {e}")

        return results


# ==================== 考试规格解析 ====================
class ExamSpecParser:
    """考试规格解析器"""

    @staticmethod
    def parse(spec_text: str) -> Dict[str, Any]:
        """
        解析考试规格文本

        示例输入: "选择题 40分 8题，简答 40分 4题，大题 20分 2题"
        """
        spec = {
            "选择题": {"score_per": 5, "count": 8, "total": 40},
            "简答题": {"score_per": 10, "count": 4, "total": 40},
            "大题": {"score_per": 10, "count": 2, "total": 20},
            "总分": 100
        }

        # 解析文本
        patterns = [
            r"([^\d]+?)\s*(\d+)\s*分\s*(\d+)\s*题",
            r"([^\d]+?)\s*(\d+)\s*题\s*(\d+)\s*分",
        ]

        for pattern in patterns:
            matches = re.findall(pattern, spec_text)
            for match in matches:
                q_type = match[0].strip()
                if "选择" in q_type:
                    key = "选择题"
                elif "简答" in q_type or "论述" in q_type:
                    key = "简答题"
                elif "大题" in q_type or "计算" in q_type or "分析" in q_type:
                    key = "大题"
                else:
                    continue

                numbers = [int(x) for x in match[1:]]
                if numbers[0] > numbers[1]:  # 40分 8题
                    spec[key] = {"score_per": numbers[0] // numbers[1], "count": numbers[1], "total": numbers[0]}
                else:  # 8题 40分
                    spec[key] = {"score_per": numbers[1] // numbers[0], "count": numbers[0], "total": numbers[1]}

        # 计算总分
        spec["总分"] = sum(v["total"] for k, v in spec.items() if isinstance(v, dict))

        return spec


# ==================== 文件下载 ====================
class PaperExporter:
    """试卷导出器"""

    @staticmethod
    def to_docx(paper: Dict, output_path: Path) -> None:
        """导出为 DOCX 格式"""
        from docx import Document
        from docx.shared import Pt, RGBColor

        doc = Document()

        # 标题
        title = doc.add_heading(paper.get("paper_title", "模拟试卷"), 0)
        title.alignment = 1  # 居中

        # 考试信息
        doc.add_paragraph(f"总分：{paper.get('total_score', 100)} 分")
        doc.add_paragraph(f"时长：{paper.get('time_limit', '120 分钟')}")
        doc.add_paragraph()

        # 写入各题型
        for section in paper.get("sections", []):
            doc.add_heading(section["section_name"], 1)

            for q in section["questions"]:
                doc.add_paragraph(f"【{q['q_id']}】{q['content']}")

                if q.get("options"):
                    for opt in q["options"]:
                        doc.add_paragraph(f"  {opt}")

                doc.add_paragraph(f"答案：{q.get('answer', '略')}")
                doc.add_paragraph(f"解析：{q.get('analysis', '略')}")
                doc.add_paragraph()

        doc.save(str(output_path))
        logger.info(f"已导出 DOCX: {output_path}")

    @staticmethod
    def to_pdf(paper: Dict, output_path: Path) -> None:
        """导出为 PDF 格式"""
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.units import cm

        doc = SimpleDocTemplate(str(output_path), pagesize=A4)
        styles = getSampleStyleSheet()
        story = []

        # 标题
        story.append(Paragraph(paper.get("paper_title", "模拟试卷"), styles["Title"]))
        story.append(Spacer(1, 0.5 * cm))

        # 考试信息
        story.append(Paragraph(f"总分：{paper.get('total_score', 100)}", styles["Normal"]))
        story.append(Spacer(1, 0.3 * cm))

        # 各题型
        for section in paper.get("sections", []):
            story.append(Paragraph(section["section_name"], styles["Heading2"]))

            for q in section["questions"]:
                story.append(Paragraph(f"【{q['q_id']}】{q['content']}", styles["Normal"]))

                if q.get("options"):
                    for opt in q["options"]:
                        story.append(Paragraph(f"  {opt}", styles["Normal"]))

                story.append(Paragraph(f"答案：{q.get('answer', '略')}", styles["Normal"]))
                story.append(Spacer(1, 0.2 * cm))

        doc.build(story)
        logger.info(f"已导出 PDF: {output_path}")

    @staticmethod
    def pack_all(papers: List[Dict], output_dir: Path, format: str = "zip") -> Path:
        """打包所有试卷"""
        import zipfile

        timestamp = time.strftime("%Y%m%d_%H%M%S")
        zip_path = output_dir / f"模拟卷_全套_{timestamp}.zip"

        with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
            for i, paper in enumerate(papers, 1):
                paper_filename = f"模拟卷_{i}_{paper.get('paper_title', f'第{i}套')}.{format}"

                if format == "docx":
                    PaperExporter.to_docx(paper, output_dir / paper_filename)
                elif format == "pdf":
                    PaperExporter.to_pdf(paper, output_dir / paper_filename)

                zf.write(output_dir / paper_filename, paper_filename)

        logger.info(f"已打包所有试卷: {zip_path}")
        return zip_path
